import os
import re
import json
import sqlite3
import uuid
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Any

from flask import (
    Flask,
    request,
    jsonify,
    render_template,
    send_from_directory
)

from werkzeug.utils import secure_filename
from dotenv import load_dotenv


# ============================================================
# OPTIONAL LIBRARIES
# ============================================================

try:
    import fitz
except Exception:
    fitz = None

try:
    from docx import Document
except Exception:
    Document = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.metrics.pairwise import cosine_similarity
except Exception:
    TfidfVectorizer = None
    cosine_similarity = None

try:
    from openai import OpenAI
except Exception:
    OpenAI = None


# ============================================================
# CONFIGURATION
# ============================================================

load_dotenv()

BASE_DIR = Path(__file__).resolve().parent

UPLOAD_DIR = BASE_DIR / "uploads"
DATA_DIR = BASE_DIR / "data"
AUDIO_DIR = BASE_DIR / "generated_audio"

UPLOAD_DIR.mkdir(exist_ok=True)
DATA_DIR.mkdir(exist_ok=True)
AUDIO_DIR.mkdir(exist_ok=True)

DATABASE = DATA_DIR / "teacher.db"

PORT = int(os.getenv("PORT", "5000"))

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "").strip()
OPENAI_MODEL = os.getenv("OPENAI_MODEL", "gpt-4o-mini")

app = Flask(__name__)

app.config["MAX_CONTENT_LENGTH"] = 25 * 1024 * 1024

ALLOWED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".txt",
    ".md"
}

client = None

if OpenAI and OPENAI_API_KEY:
    try:
        client = OpenAI(api_key=OPENAI_API_KEY)
    except Exception:
        client = None


# ============================================================
# DATABASE
# ============================================================

def get_db():
    connection = sqlite3.connect(DATABASE)
    connection.row_factory = sqlite3.Row
    return connection


def initialize_database():

    connection = get_db()

    connection.executescript(
        """
        CREATE TABLE IF NOT EXISTS learners (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL,
            level TEXT NOT NULL,
            language TEXT NOT NULL,
            teaching_style TEXT NOT NULL,
            objective TEXT,
            created_at TEXT NOT NULL
        );

        CREATE TABLE IF NOT EXISTS documents (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            learner_id INTEGER,
            filename TEXT NOT NULL,
            stored_path TEXT NOT NULL,
            content TEXT NOT NULL,
            created_at TEXT NOT NULL
        );

        CREATE TABLE IF NOT EXISTS lessons (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            learner_id INTEGER,
            topic TEXT NOT NULL,
            level TEXT,
            language TEXT,
            minutes INTEGER,
            plan_json TEXT,
            created_at TEXT NOT NULL
        );

        CREATE TABLE IF NOT EXISTS interactions (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            learner_id INTEGER,
            lesson_id INTEGER,
            concept TEXT,
            question TEXT,
            answer TEXT,
            evaluation_json TEXT,
            created_at TEXT NOT NULL
        );

        CREATE TABLE IF NOT EXISTS progress (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            learner_id INTEGER,
            lesson_id INTEGER,
            topic TEXT,
            score REAL,
            strong_areas TEXT,
            weak_areas TEXT,
            recommendation TEXT,
            created_at TEXT NOT NULL
        );
        """
    )

    connection.commit()
    connection.close()


initialize_database()


# ============================================================
# UTILITY FUNCTIONS
# ============================================================

def now():
    return datetime.utcnow().isoformat()


def clean_text(text: str) -> str:

    if not text:
        return ""

    text = text.replace("\x00", " ")

    text = re.sub(
        r"[ \t]+",
        " ",
        text
    )

    text = re.sub(
        r"\n{3,}",
        "\n\n",
        text
    )

    return text.strip()


def chunk_text(
    text: str,
    chunk_size: int = 800,
    overlap: int = 100
) -> List[str]:

    words = text.split()

    if not words:
        return []

    chunks = []

    start = 0

    while start < len(words):

        end = min(
            len(words),
            start + chunk_size
        )

        chunk = " ".join(
            words[start:end]
        )

        chunks.append(chunk)

        if end >= len(words):
            break

        start = end - overlap

    return chunks


# ============================================================
# FILE PROCESSING
# ============================================================

def extract_pdf(path: Path) -> str:

    if fitz is None:
        raise RuntimeError(
            "PyMuPDF is not installed."
        )

    document = fitz.open(path)

    pages = []

    for page in document:
        pages.append(page.get_text())

    document.close()

    return "\n\n".join(pages)


def extract_docx(path: Path) -> str:

    if Document is None:
        raise RuntimeError(
            "python-docx is not installed."
        )

    document = Document(path)

    parts = []

    for paragraph in document.paragraphs:

        text = paragraph.text.strip()

        if text:
            parts.append(text)

    for table in document.tables:

        for row in table.rows:

            values = []

            for cell in row.cells:
                values.append(cell.text.strip())

            parts.append(
                " | ".join(values)
            )

    return "\n".join(parts)


def extract_pptx(path: Path) -> str:

    if Presentation is None:
        raise RuntimeError(
            "python-pptx is not installed."
        )

    presentation = Presentation(path)

    parts = []

    for number, slide in enumerate(
        presentation.slides,
        start=1
    ):

        parts.append(
            f"[Slide {number}]"
        )

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text = shape.text.strip()

                if text:
                    parts.append(text)

    return "\n".join(parts)


def extract_document(path: Path) -> str:

    extension = path.suffix.lower()

    if extension == ".pdf":
        content = extract_pdf(path)

    elif extension == ".docx":
        content = extract_docx(path)

    elif extension == ".pptx":
        content = extract_pptx(path)

    elif extension in {".txt", ".md"}:
        content = path.read_text(
            encoding="utf-8",
            errors="ignore"
        )

    else:
        raise ValueError(
            "Unsupported file format."
        )

    return clean_text(content)


# ============================================================
# RAG RETRIEVAL
# ============================================================

def retrieve_context(
    query: str,
    document_text: str,
    top_k: int = 5
) -> List[str]:

    chunks = chunk_text(document_text)

    if not chunks:
        return []

    # --------------------------------------------------------
    # TF-IDF RETRIEVAL
    # --------------------------------------------------------

    if TfidfVectorizer and cosine_similarity:

        try:

            vectorizer = TfidfVectorizer(
                stop_words="english",
                ngram_range=(1, 2)
            )

            matrix = vectorizer.fit_transform(
                chunks
            )

            query_vector = vectorizer.transform(
                [query]
            )

            similarities = cosine_similarity(
                query_vector,
                matrix
            ).flatten()

            indexes = similarities.argsort()[::-1]

            results = []

            for index in indexes[:top_k]:

                if similarities[index] > 0:

                    results.append(
                        chunks[index]
                    )

            if results:
                return results

        except Exception:
            pass

    # --------------------------------------------------------
    # SIMPLE KEYWORD FALLBACK
    # --------------------------------------------------------

    query_words = set(
        re.findall(
            r"\w+",
            query.lower()
        )
    )

    scored = []

    for chunk in chunks:

        words = set(
            re.findall(
                r"\w+",
                chunk.lower()
            )
        )

        score = len(
            query_words.intersection(words)
        )

        scored.append(
            (score, chunk)
        )

    scored.sort(
        key=lambda x: x[0],
        reverse=True
    )

    return [
        chunk
        for score, chunk in scored[:top_k]
    ]


# ============================================================
# OPENAI
# ============================================================

def call_ai(
    system_prompt: str,
    user_prompt: str
) -> str:

    if client is None:
        return ""

    try:

        response = client.responses.create(
            model=OPENAI_MODEL,
            instructions=system_prompt,
            input=user_prompt
        )

        return response.output_text.strip()

    except Exception as error:

        print(
            "AI ERROR:",
            error
        )

        return ""


def parse_json(
    text: str,
    fallback: Any
):

    if not text:
        return fallback

    text = text.strip()

    text = re.sub(
        r"^```json\s*",
        "",
        text,
        flags=re.IGNORECASE
    )

    text = re.sub(
        r"^```\s*",
        "",
        text
    )

    text = re.sub(
        r"\s*```$",
        "",
        text
    )

    try:
        return json.loads(text)

    except Exception:
        pass

    match = re.search(
        r"(\{.*\}|\[.*\])",
        text,
        re.DOTALL
    )

    if match:

        try:
            return json.loads(
                match.group(1)
            )
        except Exception:
            pass

    return fallback


# ============================================================
# LANGUAGE
# ============================================================

def language_instruction(language):

    mapping = {

        "English":
            "Use clear simple English.",

        "Hindi":
            "Use natural Hindi. Use Devanagari when appropriate.",

        "Hinglish":
            "Use simple Hinglish in Roman script.",

        "Bengali":
            "Use natural Bengali.",

        "Tamil":
            "Use natural Tamil.",

        "Telugu":
            "Use natural Telugu."
    }

    return mapping.get(
        language,
        f"Use {language}."
    )


# ============================================================
# FALLBACK LESSON PLAN
# ============================================================

def fallback_plan(
    topic,
    level,
    minutes,
    language,
    objective
):

    if minutes <= 5:

        number_of_steps = 3

    elif minutes <= 20:

        number_of_steps = 5

    elif minutes <= 60:

        number_of_steps = 8

    else:

        number_of_steps = 10

    concepts = [

        f"Introduction to {topic}",

        f"Core concept of {topic}",

        f"Important terms in {topic}",

        f"Worked example of {topic}",

        f"Common mistakes in {topic}",

        f"Application of {topic}",

        f"Practice question for {topic}",

        f"Review of {topic}",

        f"Final understanding check",

        f"Next learning step"

    ]

    concepts = concepts[
        :number_of_steps
    ]

    steps = []

    each_time = max(
        1,
        minutes // len(concepts)
    )

    for index, concept in enumerate(
        concepts
    ):

        steps.append(
            {
                "title": concept,
                "type":
                    "explain"
                    if index < len(concepts) - 1
                    else "assessment",

                "minutes": each_time,

                "question_after": True
            }
        )

    return {

        "title":
            f"Personalized lesson: {topic}",

        "objective":
            objective or
            f"Understand the key ideas of {topic}.",

        "language":
            language,

        "level":
            level,

        "minutes":
            minutes,

        "teaching_strategy":
            "Progressive explanation with examples, questions and adaptive re-teaching.",

        "steps":
            steps,

        "visual":
            "Use diagrams, formulas, examples or subject-specific visuals.",

        "checkpoint":
            f"Explain the main idea of {topic} in your own words."
    }


# ============================================================
# AI LESSON PLANNER
# ============================================================

def generate_lesson_plan(
    topic,
    level,
    minutes,
    language,
    objective,
    context
):

    fallback = fallback_plan(
        topic,
        level,
        minutes,
        language,
        objective
    )

    if client is None:
        return fallback

    system = f"""
You are an adaptive AI teacher.

You must behave like a real teacher.

Teaching loop:

Understand
Plan
Explain
Demonstrate
Question
Evaluate
Adapt
Continue

{language_instruction(language)}

Do not behave like a simple question-answer chatbot.

Use learner level, available time, objective and source
material to design the lesson.

Return ONLY valid JSON.
"""

    user = f"""
TOPIC:
{topic}

LEARNER LEVEL:
{level}

AVAILABLE TIME:
{minutes} minutes

LEARNING OBJECTIVE:
{objective}

SOURCE MATERIAL:
{context[:14000]}

Create a personalized lesson.

The JSON must contain:

title
objective
language
level
minutes
teaching_strategy
steps
visual
checkpoint

Each step should contain:

title
type
minutes
question_after
"""

    result = parse_json(
        call_ai(
            system,
            user
        ),
        fallback
    )

    if isinstance(result, dict):
        return result

    return fallback


# ============================================================
# AI EXPLANATION
# ============================================================

def generate_explanation(
    topic,
    concept,
    level,
    language,
    teaching_style,
    context
):

    fallback = f"""
Today we are learning about {concept}.

First understand the basic idea.
Then connect it with a simple example.
After that, we will check your understanding.

Remember: the goal is not just to memorize the definition,
but to understand why the concept works.
"""

    if client is None:
        return fallback

    system = f"""
You are a patient expert teacher.

Learner level:
{level}

Teaching style:
{teaching_style}

{language_instruction(language)}

Teach progressively.

Your response should contain:

1. Intuition
2. Clear explanation
3. Example
4. Subject-aware visual suggestion
5. Short checkpoint

Do not overwhelm a beginner.
"""

    user = f"""
Topic:
{topic}

Current concept:
{concept}

Relevant source material:
{context[:10000]}
"""

    result = call_ai(
        system,
        user
    )

    return result or fallback


# ============================================================
# AI QUESTION GENERATION
# ============================================================

def generate_question(
    topic,
    concept,
    level,
    language,
    context
):

    fallback = {

        "question":
            f"What is the main idea behind {concept}?",

        "type":
            "short_answer",

        "expected":
            f"The student should explain {concept} clearly.",

        "hint":
            "Think about the definition and example."
    }

    if client is None:
        return fallback

    system = f"""
You are an expert teacher.

{language_instruction(language)}

Create one diagnostic question.

The question must test understanding,
not memorization.

Return ONLY JSON.
"""

    user = f"""
Topic:
{topic}

Concept:
{concept}

Level:
{level}

Source:
{context[:7000]}

Return:

question
type
expected
hint
"""

    result = parse_json(
        call_ai(
            system,
            user
        ),
        fallback
    )

    return (
        result
        if isinstance(result, dict)
        else fallback
    )


# ============================================================
# AI ANSWER EVALUATION
# ============================================================

def evaluate_answer(
    topic,
    concept,
    question,
    answer,
    expected,
    level,
    language
):

    fallback = {

        "correct":
            bool(answer.strip()),

        "score":
            0.7 if answer.strip()
            else 0.0,

        "misconception":
            ""
            if answer.strip()
            else
            "No answer was provided.",

        "feedback":
            "Good attempt. Let's improve the explanation.",

        "adaptation":
            "continue",

        "difficulty":
            "same",

        "re_explanation":
            "Let's connect the concept with another simple example."
    }

    if client is None:
        return fallback

    system = f"""
You are an adaptive teacher evaluator.

{language_instruction(language)}

Evaluate understanding, not just keywords.

If the learner is wrong:

1. Identify the misconception.
2. Explain why it is wrong.
3. Re-teach the concept.
4. Use another example or analogy.
5. Choose an appropriate adaptation.
6. Decide whether difficulty should change.

Return ONLY valid JSON.
"""

    user = f"""
Topic:
{topic}

Concept:
{concept}

Level:
{level}

Question:
{question}

Expected understanding:
{expected}

Student answer:
{answer}

Return:

correct
score
misconception
feedback
adaptation
difficulty
re_explanation
"""

    result = parse_json(
        call_ai(
            system,
            user
        ),
        fallback
    )

    return (
        result
        if isinstance(result, dict)
        else fallback
    )


# ============================================================
# FINAL ASSESSMENT
# ============================================================

def generate_assessment(
    topic,
    level,
    language,
    context,
    interactions
):

    fallback = {

        "questions": [

            {
                "q":
                    f"Define the central idea of {topic}.",
                "answer":
                    "A clear and correct definition."
            },

            {
                "q":
                    f"Give one example or application of {topic}.",
                "answer":
                    "A relevant example."
            },

            {
                "q":
                    f"Explain one common mistake related to {topic}.",
                "answer":
                    "A correct explanation."
            }

        ],

        "rubric":
            "Accuracy, understanding, reasoning and application."
    }

    if client is None:
        return fallback

    system = f"""
You are an expert assessment designer.

{language_instruction(language)}

Create a short final assessment.

Return ONLY JSON.
"""

    user = f"""
Topic:
{topic}

Level:
{level}

Source material:
{context[:10000]}

Previous interactions:
{json.dumps(
    interactions[-10:],
    ensure_ascii=False
)}

Create exactly 3 questions.

Return:

questions:
[
  {{
    "q": "...",
    "answer": "..."
  }}
]

rubric:
"..."
"""

    result = parse_json(
        call_ai(
            system,
            user
        ),
        fallback
    )

    return (
        result
        if isinstance(result, dict)
        else fallback
    )


# ============================================================
# ROUTES
# ============================================================

@app.get("/")
def home():

    return render_template(
        "index.html"
    )


@app.get("/health")
def health():

    return jsonify(
        {
            "status": "running",
            "ai_enabled":
                client is not None,
            "model":
                OPENAI_MODEL
        }
    )


# ============================================================
# LEARNER
# ============================================================

@app.post("/api/learner")
def create_learner():

    data = request.get_json(
        force=True
    )

    name = (
        data.get("name")
        or "Student"
    ).strip()

    level = (
        data.get("level")
        or "Beginner"
    )

    language = (
        data.get("language")
        or "English"
    )

    teaching_style = (
        data.get("teaching_style")
        or "Simple examples"
    )

    objective = (
        data.get("objective")
        or ""
    )

    connection = get_db()

    cursor = connection.execute(
        """
        INSERT INTO learners
        (
            name,
            level,
            language,
            teaching_style,
            objective,
            created_at
        )
        VALUES (?, ?, ?, ?, ?, ?)
        """,
        (
            name,
            level,
            language,
            teaching_style,
            objective,
            now()
        )
    )

    connection.commit()

    learner_id = cursor.lastrowid

    connection.close()

    return jsonify(
        {
            "success": True,
            "learner_id": learner_id
        }
    )


# ============================================================
# UPLOAD
# ============================================================

@app.post("/api/upload")
def upload_document():

    learner_id = request.form.get(
        "learner_id",
        type=int
    )

    uploaded_file = request.files.get(
        "file"
    )

    if learner_id is None:
        return jsonify(
            {"error": "Learner not found."}
        ), 400

    if uploaded_file is None:
        return jsonify(
            {"error": "Please select a file."}
        ), 400

    original_name = (
        uploaded_file.filename
        or ""
    )

    extension = Path(
        original_name
    ).suffix.lower()

    if extension not in ALLOWED_EXTENSIONS:

        return jsonify(
            {
                "error":
                    "Supported formats: PDF, DOCX, PPTX, TXT, MD."
            }
        ), 400

    safe_name = secure_filename(
        original_name
    )

    unique_name = (
        uuid.uuid4().hex[:10]
        + "_"
        + safe_name
    )

    stored_path = (
        UPLOAD_DIR
        / unique_name
    )

    uploaded_file.save(
        stored_path
    )

    try:

        content = extract_document(
            stored_path
        )

    except Exception as error:

        stored_path.unlink(
            missing_ok=True
        )

        return jsonify(
            {
                "error":
                    str(error)
            }
        ), 400

    if len(content) < 30:

        stored_path.unlink(
            missing_ok=True
        )

        return jsonify(
            {
                "error":
                    "Could not extract enough readable text."
            }
        ), 400

    connection = get_db()

    cursor = connection.execute(
        """
        INSERT INTO documents
        (
            learner_id,
            filename,
            stored_path,
            content,
            created_at
        )
        VALUES (?, ?, ?, ?, ?)
        """,
        (
            learner_id,
            original_name,
            str(stored_path),
            content,
            now()
        )
    )

    connection.commit()

    document_id = cursor.lastrowid

    connection.close()

    return jsonify(
        {
            "success": True,
            "document_id":
                document_id,
            "filename":
                original_name,
            "characters":
                len(content),
            "chunks":
                len(chunk_text(content))
        }
    )


# ============================================================
# LESSON PLAN
# ============================================================

@app.post("/api/plan")
def create_plan():

    data = request.get_json(
        force=True
    )

    learner_id = data.get(
        "learner_id"
    )

    topic = (
        data.get("topic")
        or ""
    ).strip()

    minutes = int(
        data.get(
            "minutes",
            20
        )
    )

    connection = get_db()

    learner = connection.execute(
        """
        SELECT *
        FROM learners
        WHERE id = ?
        """,
        (learner_id,)
    ).fetchone()

    documents = connection.execute(
        """
        SELECT *
        FROM documents
        WHERE learner_id = ?
        ORDER BY id DESC
        """,
        (learner_id,)
    ).fetchall()

    if learner is None:

        connection.close()

        return jsonify(
            {
                "error":
                    "Create learner first."
            }
        ), 400

    if not topic:

        if documents:

            topic = Path(
                documents[0]["filename"]
            ).stem

        else:

            topic = "General learning"

    contexts = []

    for document in documents:

        contexts.extend(
            retrieve_context(
                topic,
                document["content"],
                top_k=5
            )
        )

    context = "\n\n".join(
        contexts
    )

    plan = generate_lesson_plan(
        topic=topic,
        level=learner["level"],
        minutes=minutes,
        language=learner["language"],
        objective=learner["objective"],
        context=context
    )

    cursor = connection.execute(
        """
        INSERT INTO lessons
        (
            learner_id,
            topic,
            level,
            language,
            minutes,
            plan_json,
            created_at
        )
        VALUES (?, ?, ?, ?, ?, ?, ?)
        """,
        (
            learner_id,
            topic,
            learner["level"],
            learner["language"],
            minutes,
            json.dumps(
                plan,
                ensure_ascii=False
            ),
            now()
        )
    )

    connection.commit()

    lesson_id = cursor.lastrowid

    connection.close()

    return jsonify(
        {
            "success": True,
            "lesson_id":
                lesson_id,
            "topic":
                topic,
            "plan":
                plan
        }
    )


# ============================================================
# TEACH
# ============================================================

@app.post("/api/teach")
def teach():

    data = request.get_json(
        force=True
    )

    learner_id = data.get(
        "learner_id"
    )

    lesson_id = data.get(
        "lesson_id"
    )

    concept = (
        data.get("concept")
        or "Main concept"
    )

    connection = get_db()

    learner = connection.execute(
        """
        SELECT *
        FROM learners
        WHERE id = ?
        """,
        (learner_id,)
    ).fetchone()

    lesson = connection.execute(
        """
        SELECT *
        FROM lessons
        WHERE id = ?
        """,
        (lesson_id,)
    ).fetchone()

    documents = connection.execute(
        """
        SELECT *
        FROM documents
        WHERE learner_id = ?
        """,
        (learner_id,)
    ).fetchall()

    if learner is None or lesson is None:

        connection.close()

        return jsonify(
            {
                "error":
                    "Invalid learner or lesson."
            }
        ), 400

    contexts = []

    for document in documents:

        contexts.extend(
            retrieve_context(
                concept,
                document["content"],
                top_k=4
            )
        )

    context = "\n\n".join(
        contexts
    )

    explanation = generate_explanation(
        topic=lesson["topic"],
        concept=concept,
        level=learner["level"],
        language=learner["language"],
        teaching_style=learner["teaching_style"],
        context=context
    )

    question = generate_question(
        topic=lesson["topic"],
        concept=concept,
        level=learner["level"],
        language=learner["language"],
        context=context
    )

    connection.close()

    return jsonify(
        {
            "success": True,
            "concept":
                concept,
            "explanation":
                explanation,
            "question":
                question
        }
    )


# ============================================================
# EVALUATE
# ============================================================

@app.post("/api/evaluate")
def evaluate():

    data = request.get_json(
        force=True
    )

    learner_id = data.get(
        "learner_id"
    )

    lesson_id = data.get(
        "lesson_id"
    )

    topic = data.get(
        "topic",
        ""
    )

    concept = data.get(
        "concept",
        ""
    )

    question = data.get(
        "question",
        ""
    )

    answer = data.get(
        "answer",
        ""
    )

    expected = data.get(
        "expected",
        ""
    )

    connection = get_db()

    learner = connection.execute(
        """
        SELECT *
        FROM learners
        WHERE id = ?
        """,
        (learner_id,)
    ).fetchone()

    if learner is None:

        connection.close()

        return jsonify(
            {
                "error":
                    "Learner not found."
            }
        ), 400

    evaluation = evaluate_answer(
        topic=topic,
        concept=concept,
        question=question,
        answer=answer,
        expected=expected,
        level=learner["level"],
        language=learner["language"]
    )

    connection.execute(
        """
        INSERT INTO interactions
        (
            learner_id,
            lesson_id,
            concept,
            question,
            answer,
            evaluation_json,
            created_at
        )
        VALUES (?, ?, ?, ?, ?, ?, ?)
        """,
        (
            learner_id,
            lesson_id,
            concept,
            question,
            answer,
            json.dumps(
                evaluation,
                ensure_ascii=False
            ),
            now()
        )
    )

    connection.commit()
    connection.close()

    return jsonify(
        {
            "success": True,
            **evaluation
        }
    )


# ============================================================
# FINAL ASSESSMENT
# ============================================================

@app.post("/api/assessment")
def assessment():

    data = request.get_json(
        force=True
    )

    learner_id = data.get(
        "learner_id"
    )

    lesson_id = data.get(
        "lesson_id"
    )

    connection = get_db()

    learner = connection.execute(
        """
        SELECT *
        FROM learners
        WHERE id = ?
        """,
        (learner_id,)
    ).fetchone()

    lesson = connection.execute(
        """
        SELECT *
        FROM lessons
        WHERE id = ?
        """,
        (lesson_id,)
    ).fetchone()

    documents = connection.execute(
        """
        SELECT *
        FROM documents
        WHERE learner_id = ?
        """,
        (learner_id,)
    ).fetchall()

    interactions = connection.execute(
        """
        SELECT *
        FROM interactions
        WHERE lesson_id = ?
        ORDER BY id
        """,
        (lesson_id,)
    ).fetchall()

    if learner is None or lesson is None:

        connection.close()

        return jsonify(
            {
                "error":
                    "Invalid lesson."
            }
        ), 400

    context = "\n\n".join(
        document["content"][:5000]
        for document in documents
    )

    interaction_data = []

    for interaction in interactions:

        interaction_data.append(
            dict(interaction)
        )

    result = generate_assessment(
        topic=lesson["topic"],
        level=learner["level"],
        language=learner["language"],
        context=context,
        interactions=interaction_data
    )

    connection.close()

    return jsonify(
        {
            "success": True,
            **result
        }
    )


# ============================================================
# SAVE REPORT
# ============================================================

@app.post("/api/report")
def save_report():

    data = request.get_json(
        force=True
    )

    learner_id = data.get(
        "learner_id"
    )

    lesson_id = data.get(
        "lesson_id"
    )

    topic = data.get(
        "topic",
        ""
    )

    score = float(
        data.get(
            "score",
            0
        )
    )

    strong_areas = data.get(
        "strong",
        []
    )

    weak_areas = data.get(
        "weak",
        []
    )

    recommendation = data.get(
        "recommendation",
        ""
    )

    connection = get_db()

    connection.execute(
        """
        INSERT INTO progress
        (
            learner_id,
            lesson_id,
            topic,
            score,
            strong_areas,
            weak_areas,
            recommendation,
            created_at
        )
        VALUES (?, ?, ?, ?, ?, ?, ?, ?)
        """,
        (
            learner_id,
            lesson_id,
            topic,
            score,
            json.dumps(
                strong_areas,
                ensure_ascii=False
            ),
            json.dumps(
                weak_areas,
                ensure_ascii=False
            ),
            recommendation,
            now()
        )
    )

    connection.commit()

    history = connection.execute(
        """
        SELECT *
        FROM progress
        WHERE learner_id = ?
        ORDER BY id DESC
        """,
        (learner_id,)
    ).fetchall()

    connection.close()

    return jsonify(
        {
            "success": True,
            "history":
                [
                    dict(row)
                    for row in history
                ]
        }
    )


# ============================================================
# PROFILE
# ============================================================

@app.get("/api/profile/<int:learner_id>")
def profile(learner_id):

    connection = get_db()

    learner = connection.execute(
        """
        SELECT *
        FROM learners
        WHERE id = ?
        """,
        (learner_id,)
    ).fetchone()

    lessons = connection.execute(
        """
        SELECT *
        FROM lessons
        WHERE learner_id = ?
        ORDER BY id DESC
        """,
        (learner_id,)
    ).fetchall()

    progress = connection.execute(
        """
        SELECT *
        FROM progress
        WHERE learner_id = ?
        ORDER BY id DESC
        """,
        (learner_id,)
    ).fetchall()

    connection.close()

    return jsonify(
        {
            "learner":
                dict(learner)
                if learner
                else None,

            "lessons":
                [
                    dict(row)
                    for row in lessons
                ],

            "progress":
                [
                    dict(row)
                    for row in progress
                ]
        }
    )


# ============================================================
# SERVER
# ============================================================

if __name__ == "__main__":

    print()
    print("=" * 60)
    print("        AI TEACHER - HACKATHON PROJECT")
    print("=" * 60)
    print(
        "AI enabled:",
        bool(client)
    )
    print(
        "Model:",
        OPENAI_MODEL
    )
    print(
        "Open:",
        f"http://127.0.0.1:{PORT}"
    )
    print("=" * 60)
    print()

    app.run(
        host="127.0.0.1",
        port=PORT,
        debug=True
    )